# Build a PowerPoint (.pptx) architecture deck for the Retail Banking Assistant use case.
# Native shapes/tables => fully editable after importing into Google Slides.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- palette ----
INK   = RGBColor(0x0F,0x1E,0x2E)
MUTED = RGBColor(0x5B,0x6B,0x7B)
BG    = RGBColor(0xEE,0xF2,0xF6)
CARD  = RGBColor(0xFF,0xFF,0xFF)
LINE  = RGBColor(0xD8,0xE0,0xE8)
NAVY  = RGBColor(0x0B,0x20,0x38)
NAVY2 = RGBColor(0x12,0x3A,0x5E)
WHITE = RGBColor(0xFF,0xFF,0xFF)
ACCENT= RGBColor(0x0E,0xA5,0xA4)
CLIENT= RGBColor(0x16,0xA3,0x4A)
ORCH  = RGBColor(0x63,0x66,0xF1)
MCP   = RGBColor(0x0E,0xA5,0xA4)
A2A   = RGBColor(0xF5,0x9E,0x0B)
DATA  = RGBColor(0x64,0x74,0x8B)
EXT   = RGBColor(0xA8,0x55,0xF7)
LIGHTTEAL = RGBColor(0xEA,0xF6,0xF6)
CARDDK = RGBColor(0x0E,0x2C,0x46)

FONT = "Segoe UI"
EMU = 914400
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s

def _set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=FONT, space_after=6):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    runs = text if isinstance(text, list) else [(text, {})]
    for i,(t,opt) in enumerate(runs):
        r = p.add_run(); r.text = t
        f = r.font; f.size = Pt(opt.get("size",size)); f.name = font
        f.bold = opt.get("bold",bold); f.color.rgb = opt.get("color",color)
    return p

def textbox(s, x,y,w,h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font=FONT, space_after=6):
    tb = s.shapes.add_textbox(x,y,w,h); tf = tb.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    _set_text(tf, text, size, color, bold, align, font, space_after)
    return tb

def bullets(s, x,y,w,h, items, size=17, color=INK, gap=10, dot=ACCENT, font=FONT):
    tb = s.shapes.add_textbox(x,y,w,h); tf = tb.text_frame; tf.word_wrap=True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first=True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.space_after = Pt(gap)
        # bullet char
        rb = p.add_run(); rb.text="■  "; rb.font.size=Pt(11); rb.font.color.rgb=dot; rb.font.name=font
        parts = it if isinstance(it,list) else [(it,{})]
        for t,opt in parts:
            r=p.add_run(); r.text=t; f=r.font
            f.size=Pt(opt.get("size",size)); f.name=font; f.bold=opt.get("bold",False)
            f.color.rgb=opt.get("color",color)
    return tb

def box(s, x,y,w,h, fill, radius=0.10, line=None, shadow=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,h)
    try: shp.adjustments[0]=radius
    except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(1)
    shp.shadow.inherit=False
    if shadow:
        el = shp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'), {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        ef.append(sh)
        clr = sh.makeelement(qn('a:srgbClr'), {'val':'0F1E2E'}); sh.append(clr)
        a = clr.makeelement(qn('a:alpha'), {'val':'18000'}); clr.append(a)
    return shp

def boxtext(s, shp, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(6); tf.margin_right=Pt(6); tf.margin_top=Pt(4); tf.margin_bottom=Pt(4)
    first=True
    for (t,size,color,bold) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first=False; p.alignment=align
        r=p.add_run(); r.text=t; f=r.font
        f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.name=FONT

def eyebrow(s, text, y=Inches(0.55), color=ACCENT):
    textbox(s, Inches(0.8), y, Inches(11.7), Inches(0.4),
            text.upper(), size=13, color=color, bold=True)

def title(s, text, y=Inches(0.95), color=INK, size=32):
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y+Inches(0.12), Inches(0.5), Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background(); bar.shadow.inherit=False
    textbox(s, Inches(1.45), y, Inches(11.1), Inches(0.7), text, size=size, color=color, bold=True)

def footer(s, right="Retail Banking Assistant", dark=False):
    c = RGBColor(0x9F,0xB6,0xC6) if dark else MUTED
    textbox(s, Inches(0.8), Inches(7.02), Inches(4), Inches(0.35),
            [("FLOGO", {"bold":True,"color":c}), ("·", {"bold":True,"color":ACCENT}), ("AI", {"bold":True,"color":c})],
            size=13)
    textbox(s, Inches(8.5), Inches(7.02), Inches(4.0), Inches(0.35), right, size=12, color=c, align=PP_ALIGN.RIGHT)

def chip(s, x, y, text, fill, w=None):
    w = w or Inches(0.02*len(text)+0.5)
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,Inches(0.42))
    try: c.adjustments[0]=0.5
    except Exception: pass
    c.fill.solid(); c.fill.fore_color.rgb=fill; c.line.fill.background(); c.shadow.inherit=False
    boxtext(s, c, [(text,13,WHITE,True)])
    return c

def connector(s, x1,y1,x2,y2, color=RGBColor(0x41,0x58,0x6B), width=2.2, dashed=False, arrow=True):
    # OOXML requires integer EMU coordinates and NON-ZERO extents (cx>0, cy>0);
    # Google Slides rejects zero/float extents as "corrupt". Coerce + nudge.
    x1,y1,x2,y2 = int(round(x1)),int(round(y1)),int(round(x2)),int(round(y2))
    if abs(x2-x1) < 9525: x2 = x1 + 9525   # ~0.01" — invisible, keeps cx>0
    if abs(y2-y1) < 9525: y2 = y1 + 9525   # keeps cy>0
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1,y1,x2,y2)
    cn.line.color.rgb=color; cn.line.width=Pt(width); cn.shadow.inherit=False
    ln = cn.line._get_or_add_ln()
    if dashed:
        d = ln.makeelement(qn('a:prstDash'), {'val':'dash'}); ln.append(d)
    if arrow:
        he = ln.makeelement(qn('a:tailEnd'), {'type':'triangle','w':'med','len':'med'}); ln.append(he)
    return cn

def card(s, x,y,w,h, topcolor=ACCENT):
    box(s, x,y,w,h, CARD, radius=0.06, line=LINE)
    top = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,Inches(0.07))
    top.fill.solid(); top.fill.fore_color.rgb=topcolor; top.line.fill.background(); top.shadow.inherit=False

def table(s, x,y,w, rows, colw, header=True, fontsize=13.5, rowh=Inches(0.42)):
    nrows=len(rows); ncols=len(rows[0])
    gt = s.shapes.add_table(nrows,ncols,x,y,w,rowh*nrows).table
    gt.first_row=header; gt.horz_banding=False
    for ci,cw in enumerate(colw): gt.columns[ci].width=cw
    for ri,row in enumerate(rows):
        gt.rows[ri].height=rowh
        for ci,val in enumerate(row):
            cell=gt.cell(ri,ci)
            cell.margin_left=Pt(8); cell.margin_right=Pt(6); cell.margin_top=Pt(3); cell.margin_bottom=Pt(3)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            if header and ri==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
                col=WHITE; bold=True; sz=12
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb=CARD
                col=INK; bold=False; sz=fontsize
            tf=cell.text_frame; tf.word_wrap=True
            p=tf.paragraphs[0]; r=p.add_run(); r.text=str(val)
            f=r.font; f.size=Pt(sz); f.name=FONT; f.bold=bold; f.color.rgb=col
    return gt

# =========================================================================
# 1 TITLE
s = slide(NAVY)
# subtle top band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, SW, Inches(0.14))
band.fill.solid(); band.fill.fore_color.rgb=ACCENT; band.line.fill.background(); band.shadow.inherit=False
textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
        "AGENTIC AI USE CASE · TIBCO FLOGO ENTERPRISE", size=15, color=RGBColor(0x7F,0xD4,0xD1), bold=True)
textbox(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(2.0),
        "Retail Banking\nSelf-Service Assistant", size=52, color=WHITE, bold=True)
textbox(s, Inches(0.9), Inches(4.35), Inches(10.8), Inches(1.2),
        "A conversational AI that lets a bank customer check balances, transactions, cards & loans — "
        "and safely dispute a charge or block a card — all in natural language.",
        size=20, color=RGBColor(0xC7,0xD7,0xE3))
for i,(t,c) in enumerate([("AI Orchestrator",ORCH),("MCP · 7 read tools",MCP),("A2A · 3 write agents",A2A),("PostgreSQL",DATA)]):
    chip(s, Inches(0.9+i*3.0), Inches(5.7), t, c, w=Inches(2.8))
footer(s, "Architecture Overview", dark=True)

# =========================================================================
# 2 PROBLEM / VALUE
s = slide(); eyebrow(s,"Why"); title(s,"The problem we automate")
data=[("Today",DATA,"Call-center bound","Routine questions and low-risk actions (dispute, block card) tie up agents and leave customers on hold."),
      ("With this assistant",ORCH,"Self-service in chat","The customer asks in plain English. The AI answers from live data and performs confirmed actions instantly, 24/7."),
      ("Outcome",ACCENT,"Faster + safer","Reads are instant & stateless; every write is confirmed, bounded, and auditable. Risky requests are declined.")]
cw=Inches(3.9); gap=Inches(0.25); x0=Inches(0.8); y0=Inches(2.0); h=Inches(3.6)
for i,(k,tc,hd,body) in enumerate(data):
    x=x0+i*(cw+gap); card(s,x,y0,cw,h,tc)
    textbox(s,x+Inches(0.28),y0+Inches(0.32),cw-Inches(0.5),Inches(0.35),k.upper(),size=12,color=MUTED,bold=True)
    textbox(s,x+Inches(0.28),y0+Inches(0.72),cw-Inches(0.5),Inches(0.6),hd,size=21,color=INK,bold=True)
    textbox(s,x+Inches(0.28),y0+Inches(1.5),cw-Inches(0.5),Inches(1.9),body,size=16,color=MUTED,space_after=0)
footer(s)

# =========================================================================
# 3 CAPABILITIES
s = slide(); eyebrow(s,"What the customer can do"); title(s,"Capabilities")
cw=Inches(5.9); h=Inches(4.2); y0=Inches(2.0)
card(s,Inches(0.8),y0,cw,h,MCP)
chip(s,Inches(1.1),y0+Inches(0.3),"Ask & look up  (read-only)",MCP,w=Inches(3.4))
bullets(s,Inches(1.1),y0+Inches(1.0),cw-Inches(0.6),Inches(3.0),
        ["Account balances & details","Recent transactions & a specific charge","Debit / credit card details",
         "Loan balance, EMI, rate, next due date","Branch / ATM locations & hours","Status of an existing dispute"], size=16, dot=MCP)
card(s,Inches(6.9),y0,cw,h,A2A)
chip(s,Inches(7.2),y0+Inches(0.3),"Take action  (state-changing)",A2A,w=Inches(3.6))
bullets(s,Inches(7.2),y0+Inches(1.0),cw-Inches(0.6),Inches(1.7),
        [[("Dispute a transaction",{"bold":True}),(" they don’t recognize",{})],
         [("Block",{"bold":True}),(" a lost / stolen / compromised card",{})],
         [("Email",{"bold":True}),(" a confirmation of what was done",{})]], size=16, dot=A2A)
textbox(s,Inches(7.2),y0+Inches(3.0),cw-Inches(0.6),Inches(1.0),
        [("Out of scope (declined): ",{"bold":True,"color":INK}),
         ("loan approvals, investment/tax advice, opening/closing accounts, external wires.",{"color":MUTED})], size=15)
footer(s)

# =========================================================================
# 4 ARCHITECTURE DIAGRAM
s = slide(); eyebrow(s,"How it fits together"); title(s,"Architecture — the 3-app pattern")
# boxes
def cbox(x,y,w,h,fill,lines):
    b=box(s,x,y,w,h,fill,radius=0.12); boxtext(s,b,lines); return b
# Chatbot
cbox(Inches(4.7),Inches(1.7),Inches(3.9),Inches(0.8),CLIENT,
     [("Agentic Chatbot UI",15,WHITE,True),("browser · demos/Agentic_AI/Chatbot",11,RGBColor(0xEA,0xF7,0xEF),False)])
# Orchestrator
cbox(Inches(4.6),Inches(3.0),Inches(4.1),Inches(1.0),ORCH,
     [("BankingAIOrchestrator",16,WHITE,True),("#wsserver → AI Agent → reply",11.5,RGBColor(0xE7,0xE7,0xFF),False),
      ("intent routing · confirm before writes",11.5,RGBColor(0xE7,0xE7,0xFF),False)])
# MCP
cbox(Inches(0.9),Inches(4.7),Inches(3.5),Inches(1.05),MCP,
     [("BankingMCPServer",15,WHITE,True),("7 read-only tools",12,RGBColor(0xE6,0xFA,0xFA),False),
      ("accounts · txns · cards · loans …",11.5,RGBColor(0xE6,0xFA,0xFA),False)])
# A2A
cbox(Inches(8.9),Inches(4.7),Inches(3.5),Inches(1.05),A2A,
     [("BankingA2AServers",15,WHITE,True),("dispute · block-card · email",12,RGBColor(0xFF,0xF3,0xDE),False),
      ("write workflows",11.5,RGBColor(0xFF,0xF3,0xDE),False)])
# Postgres
cbox(Inches(4.6),Inches(6.25),Inches(4.1),Inches(0.7),DATA,
     [("PostgreSQL — database “banking” · 7 tables",13.5,WHITE,True)])
# LLM + SMTP
cbox(Inches(11.3),Inches(2.95),Inches(1.7),Inches(0.75),EXT,[("OpenAI LLM",12.5,WHITE,True),("gpt-5.5",11,RGBColor(0xF3,0xE8,0xFF),False)])
cbox(Inches(11.3),Inches(4.85),Inches(1.7),Inches(0.7),EXT,[("SMTP",12.5,WHITE,True),("Gmail SSL:465",10.5,RGBColor(0xF3,0xE8,0xFF),False)])
# connectors
connector(s, Inches(6.55),Inches(2.5), Inches(6.55),Inches(3.0))
connector(s, Inches(5.6),Inches(4.0), Inches(3.1),Inches(4.7))
connector(s, Inches(7.7),Inches(4.0), Inches(10.2),Inches(4.7))
connector(s, Inches(2.6),Inches(5.75), Inches(5.2),Inches(6.25))
connector(s, Inches(10.6),Inches(5.75), Inches(8.0),Inches(6.25))
connector(s, Inches(10.6),Inches(5.1), Inches(11.3),Inches(5.15))            # A2A->SMTP
connector(s, Inches(8.7),Inches(3.35), Inches(11.3),Inches(3.25), color=EXT, dashed=True)  # orch->LLM
connector(s, Inches(12.1),Inches(3.7), Inches(11.9),Inches(4.85), color=EXT, dashed=True)  # LLM->A2A(ish)
# edge labels
textbox(s, Inches(6.7),Inches(2.55),Inches(4.2),Inches(0.3),"WebSocket · ws://localhost:8088/banking",size=11,color=RGBColor(0x33,0x45,0x5A),bold=True)
textbox(s, Inches(1.4),Inches(4.12),Inches(2.4),Inches(0.3),"MCP (HTTP) :9096",size=11,color=RGBColor(0x33,0x45,0x5A),bold=True)
textbox(s, Inches(8.0),Inches(4.12),Inches(2.6),Inches(0.3),"A2A :8710-8712",size=11,color=RGBColor(0x33,0x45,0x5A),bold=True,align=PP_ALIGN.RIGHT)
footer(s)

# =========================================================================
# 5 WHY 3 APPS
s = slide(); eyebrow(s,"Design principle"); title(s,"Why three separate apps?")
data=[("MCP Server",MCP,"Reads only","Stateless lookups, safe to retry. The LLM picks a tool by intent and filters rows. No side effects."),
      ("A2A Servers",A2A,"Writes only","Each state-changing action is its own agent with its own guardrails, port, and deploy/scale lifecycle."),
      ("Orchestrator",ORCH,"The brain","WebSocket chat. Classifies intent, calls read tools or hands off to write agents, keeps memory, enforces scope.")]
cw=Inches(3.9); gap=Inches(0.25); x0=Inches(0.8); y0=Inches(2.0); h=Inches(3.0)
for i,(k,tc,hd,body) in enumerate(data):
    x=x0+i*(cw+gap); card(s,x,y0,cw,h,tc)
    chip(s,x+Inches(0.28),y0+Inches(0.3),k,tc,w=Inches(2.1))
    textbox(s,x+Inches(0.28),y0+Inches(0.95),cw-Inches(0.5),Inches(0.5),hd,size=20,color=INK,bold=True)
    textbox(s,x+Inches(0.28),y0+Inches(1.6),cw-Inches(0.5),Inches(1.3),body,size=15,color=MUTED,space_after=0)
textbox(s,Inches(0.8),Inches(5.4),Inches(11.7),Inches(1.0),
        [("Separation of read vs. write is the core safety property: ",{"color":INK}),
         ("reads can’t mutate; writes are explicit, confirmed, and independently governable.",{"bold":True,"color":INK})],
        size=18)
footer(s)

# =========================================================================
# 6 MCP TOOLS
s = slide(); eyebrow(s,"BankingMCPServer · port 9096 · /retail-banking")
title(s,"MCP Server — 7 read tools")
rows=[["Tool","Returns","Table"],
      ["GetCustomerProfile","name, phone, email","customers"],
      ["GetAccounts","type, balance, currency, status","accounts"],
      ["GetTransactions","date, merchant, amount, type, status","transactions"],
      ["GetCards","masked no., type, network, status, expiry","cards"],
      ["GetLoans","type, outstanding, rate, EMI, next due","loans"],
      ["GetDisputes","txn, reason, status, est. resolution","disputes"],
      ["GetBranches","name, address, city, phone, hours","branches"]]
table(s, Inches(0.8), Inches(1.95), Inches(11.7), rows, [Inches(3.2),Inches(5.6),Inches(2.9)], rowh=Inches(0.46))
textbox(s,Inches(0.8),Inches(6.2),Inches(11.7),Inches(0.6),
        [("Each tool is a tiny flow:  ",{"color":MUTED}),
         ("#noop → #query (SELECT *) → #actreturn",{"bold":True,"color":RGBColor(0x0B,0x6B,0x6A)}),
         ("   — rich descriptions let the LLM choose correctly.",{"color":MUTED})], size=15)
footer(s)

# =========================================================================
# 7 A2A AGENTS
s = slide(); eyebrow(s,"BankingA2AServers · ports 8710 · 8711 · 8712")
title(s,"A2A Servers — 3 write agents")
rows=[["Agent","Port","Workflow"],
      ["dispute_transaction_agent","8710","validate the transaction (#query) → INSERT a dispute (OPEN, +7 business days)"],
      ["block_card_agent","8711","UPDATE the card → status = BLOCKED"],
      ["send_confirmation_email","8712","#sendmail → one confirmation email (Gmail SSL:465)"]]
table(s, Inches(0.8), Inches(1.95), Inches(11.7), rows, [Inches(3.4),Inches(1.1),Inches(7.2)], rowh=Inches(0.62))
gd=[("Guardrail","Confirm before every write"),("Bounded","Max 3 attempts per action"),("Email rule","Sent once, last, only if asked")]
cw=Inches(3.7); gap=Inches(0.3); x0=Inches(0.8); y0=Inches(4.6); h=Inches(1.4)
for i,(k,v) in enumerate(gd):
    x=x0+i*(cw+gap); card(s,x,y0,cw,h,ACCENT)
    textbox(s,x+Inches(0.28),y0+Inches(0.28),cw-Inches(0.5),Inches(0.35),k.upper(),size=12,color=MUTED,bold=True)
    textbox(s,x+Inches(0.28),y0+Inches(0.68),cw-Inches(0.5),Inches(0.6),v,size=16,color=INK)
footer(s)

# =========================================================================
# 8 ORCHESTRATOR
s = slide(); eyebrow(s,"BankingAIOrchestrator · WebSocket 8088 · /banking")
title(s,"Orchestrator — intent routing & guardrails")
cw=Inches(5.9); h=Inches(3.5); y0=Inches(2.0)
card(s,Inches(0.8),y0,cw,h,ORCH)
textbox(s,Inches(1.1),y0+Inches(0.28),cw-Inches(0.6),Inches(0.35),"ROUTING LOGIC (SYSTEM PROMPT)",size=12,color=MUTED,bold=True)
bullets(s,Inches(1.1),y0+Inches(0.8),cw-Inches(0.6),Inches(2.6),
        [[("Identify customer by ",{}),("CUST-2026-NNNNN",{"bold":True})],
         [("Balance / cards / loans / branch → ",{}),("MCP tool",{"bold":True})],
         [("“Don’t recognize this charge” → ",{}),("dispute agent",{"bold":True})],
         [("“Lost / stolen card” → ",{}),("block-card agent",{"bold":True})],
         [("“Email me” → ",{}),("email agent (once, last)",{"bold":True})]], size=15.5, dot=ORCH)
card(s,Inches(6.9),y0,cw,h,ACCENT)
textbox(s,Inches(7.2),y0+Inches(0.28),cw-Inches(0.6),Inches(0.35),"SAFETY & BEHAVIOR",size=12,color=MUTED,bold=True)
bullets(s,Inches(7.2),y0+Inches(0.8),cw-Inches(0.6),Inches(2.6),
        ["Confirm before any write","Reveal only the requesting customer’s data",
         "Session memory for multi-turn chat","Decline out-of-scope requests politely","Report partial results on failure"], size=15.5)
textbox(s,Inches(0.8),Inches(5.75),Inches(11.7),Inches(0.6),
        [("Flow:  ",{"color":MUTED}),("#wsserver → #agentactivity (MCP + 3 A2A wired in) → #wswritedata",{"bold":True,"color":RGBColor(0x0B,0x6B,0x6A)})],size=15)
footer(s)

# =========================================================================
# 9 DATA MODEL
s = slide(); eyebrow(s,"PostgreSQL · database “banking”"); title(s,"Data model — 7 tables")
rows=[["Table","Purpose","Accessed by"],
      ["customers","master record, keyed by customer_id","profile · email recipient"],
      ["accounts","checking / savings / credit balances","GetAccounts"],
      ["transactions","history + the disputable charge","GetTransactions · dispute"],
      ["cards","debit / credit cards","GetCards · block (UPDATE)"],
      ["loans","home / auto / personal loans","GetLoans"],
      ["disputes","filed disputes","GetDisputes · file (INSERT)"],
      ["branches","branch / ATM reference","GetBranches"]]
table(s, Inches(0.8), Inches(1.9), Inches(11.7), rows, [Inches(2.8),Inches(5.3),Inches(3.6)], rowh=Inches(0.44))
stats=[("7","customers"),("15","transactions"),("7","cards"),("1","pre-seeded dispute")]
x0=Inches(0.9); y0=Inches(6.05)
for i,(n,l) in enumerate(stats):
    textbox(s,x0+i*Inches(3.0),y0,Inches(2.8),Inches(0.6),n,size=40,color=ACCENT,bold=True)
    textbox(s,x0+i*Inches(3.0),y0+Inches(0.7),Inches(2.8),Inches(0.3),l,size=14,color=MUTED)
footer(s)

# =========================================================================
# 10 SEQUENCE (dark)
s = slide(NAVY); eyebrow(s,"End-to-end",color=RGBColor(0x7F,0xD4,0xD1))
textbox(s,Inches(0.8),Inches(0.95),Inches(11.9),Inches(0.7),
        "Walkthrough — “Dispute my $249.99 charge, then email me”",size=28,color=WHITE,bold=True)
steps=[("1","Customer asks","“I don’t recognize a $249.99 QUICKPAY charge”"),
       ("2","Orchestrator","GetTransactions → finds TXN-50003, asks to confirm"),
       ("3","Customer confirms","“Yes, dispute it”"),
       ("4","dispute agent","INSERT dispute DSP-… · status OPEN · +7 days"),
       ("5","email agent","sends ONE confirmation (only because asked)"),
       ("6","Reply streamed","summary back over the WebSocket to the chat")]
bw=Inches(3.85); bh=Inches(1.5); gx=Inches(0.28); gy=Inches(0.35)
x0=Inches(0.8); y0=Inches(2.0)
pos=[]
half_w=int(bw/2); half_h=int(bh/2)
for i,(n,hd,body) in enumerate(steps):
    # snake layout: top row 1-2-3 left->right, bottom row 4-5-6 right->left
    if i<3: col=i; row=0
    else:   col=2-(i-3); row=1
    x=x0+col*(bw+gx); y=y0+row*(bh+gy)
    pos.append((x,y))
    b=box(s,x,y,bw,bh,CARDDK,radius=0.08,line=NAVY2)
    textbox(s,x+Inches(0.25),y+Inches(0.18),Inches(0.6),Inches(0.4),n,size=20,color=RGBColor(0x7F,0xD4,0xD1),bold=True)
    textbox(s,x+Inches(0.75),y+Inches(0.2),bw-Inches(1.0),Inches(0.4),hd,size=15,color=WHITE,bold=True)
    textbox(s,x+Inches(0.25),y+Inches(0.72),bw-Inches(0.5),Inches(0.7),body,size=12.5,color=RGBColor(0xBC,0xD4,0xE2),space_after=0)
# flow arrows (snake): 1->2->3 (top, L->R), 3 down to 4, 4->5->6 (bottom, R->L)
tl=RGBColor(0x7F,0xD4,0xD1)
connector(s, pos[0][0]+bw, pos[0][1]+half_h, pos[1][0], pos[1][1]+half_h, color=tl, width=2)
connector(s, pos[1][0]+bw, pos[1][1]+half_h, pos[2][0], pos[2][1]+half_h, color=tl, width=2)
connector(s, pos[2][0]+half_w, pos[2][1]+bh, pos[3][0]+half_w, pos[3][1], color=tl, width=2)
connector(s, pos[3][0], pos[3][1]+half_h, pos[4][0]+bw, pos[4][1]+half_h, color=tl, width=2)
connector(s, pos[4][0], pos[4][1]+half_h, pos[5][0]+bw, pos[5][1]+half_h, color=tl, width=2)
textbox(s,Inches(0.8),Inches(6.35),Inches(11.9),Inches(0.7),
        [("Reads gather context → the customer confirms → a single write fires → one email → streamed reply.  ",{"color":RGBColor(0xC7,0xD7,0xE3)}),
         ("Nothing is written without confirmation.",{"bold":True,"color":WHITE})],size=15)
footer(s, dark=True)

# =========================================================================
# 11 PORTS / DEPLOY
s = slide(); eyebrow(s,"Deploy & run"); title(s,"Ports & start order")
cw=Inches(5.9); h=Inches(4.3); y0=Inches(2.0)
card(s,Inches(0.8),y0,cw,h,MCP)
textbox(s,Inches(1.1),y0+Inches(0.28),cw-Inches(0.6),Inches(0.35),"ENDPOINTS",size=12,color=MUTED,bold=True)
rows=[["App","Port / path"],
      ["MCP Server","9096 · /retail-banking"],
      ["A2A — dispute","8710"],
      ["A2A — block card","8711"],
      ["A2A — email","8712"],
      ["Orchestrator (WS)","8088 · /banking"]]
table(s, Inches(1.1), y0+Inches(0.75), cw-Inches(0.6), rows, [Inches(2.7),Inches(2.5)], rowh=Inches(0.5))
card(s,Inches(6.9),y0,cw,h,ORCH)
textbox(s,Inches(7.2),y0+Inches(0.28),cw-Inches(0.6),Inches(0.35),"START ORDER",size=12,color=MUTED,bold=True)
bullets(s,Inches(7.2),y0+Inches(0.8),cw-Inches(0.6),Inches(2.6),
        [[("1.  Load DB — ",{"bold":True}),("database.sql, then reset_data.sql",{})],
         [("2.  Start MCP · 9096",{"bold":True})],
         [("3.  Start A2A · 8710-8712",{"bold":True})],
         [("4.  Start Orchestrator · 8088",{"bold":True})],
         [("5.  Connect chatbot → ",{"bold":True}),("ws://localhost:8088/banking",{})]], size=15.5, dot=ORCH)
textbox(s,Inches(7.2),y0+Inches(3.55),cw-Inches(0.6),Inches(0.6),
        "Reuses existing OpenAI / PostgreSQL / SMTP connections; only the DB name changes to “banking”.",
        size=13.5,color=MUTED)
footer(s)

# =========================================================================
# 12 CLOSING (dark)
s = slide(NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, SW, Inches(0.14))
band.fill.solid(); band.fill.fore_color.rgb=ACCENT; band.line.fill.background(); band.shadow.inherit=False
textbox(s,Inches(0.9),Inches(1.2),Inches(11),Inches(0.4),"SUMMARY",size=14,color=RGBColor(0x7F,0xD4,0xD1),bold=True)
textbox(s,Inches(0.9),Inches(1.7),Inches(11.6),Inches(1.0),"One pattern, any vertical",size=42,color=WHITE,bold=True)
textbox(s,Inches(0.9),Inches(2.8),Inches(11.2),Inches(1.0),
        "Reads as MCP tools, writes as A2A agents, a WebSocket orchestrator as the brain — the same shape "
        "powers Hospital, Telecom, and now Retail Banking.",size=19,color=RGBColor(0xC7,0xD7,0xE3))
cards=[("Safe","Reads can’t mutate; writes are confirmed, bounded, auditable."),
       ("Modular","Each capability deploys & scales independently."),
       ("Extensible","Add a tool or an agent without touching the rest.")]
cw=Inches(3.9); gap=Inches(0.25); x0=Inches(0.9); y0=Inches(4.2); h=Inches(1.7)
for i,(k,v) in enumerate(cards):
    x=x0+i*(cw+gap); b=box(s,x,y0,cw,h,CARDDK,radius=0.06,line=NAVY2)
    textbox(s,x+Inches(0.28),y0+Inches(0.25),cw-Inches(0.5),Inches(0.35),k.upper(),size=13,color=RGBColor(0x7F,0xD4,0xD1),bold=True)
    textbox(s,x+Inches(0.28),y0+Inches(0.7),cw-Inches(0.5),Inches(0.9),v,size=15,color=RGBColor(0xD7,0xE6,0xEF),space_after=0)
textbox(s,Inches(0.9),Inches(6.35),Inches(11.5),Inches(0.4),
        "demos/Agentic_AI/Retail_Banking_Assistant_Use_Case",size=14,color=RGBColor(0x9F,0xD8,0xD5))
footer(s,"Thank you",dark=True)

out = "RetailBankingAssistant_Architecture.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
