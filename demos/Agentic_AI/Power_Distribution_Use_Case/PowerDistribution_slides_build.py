# -*- coding: utf-8 -*-
"""
Build a concise 6-slide Agentic AI demo deck for the Electric Power Distribution
use case, ON the official TIBCO 2026 corporate template (logo/master/theme carry over).
No customer name anywhere. Run with:  PYTHONIOENCODING=utf-8 python PowerDistribution_slides_build.py
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE = r"C:/Users/nshah/Downloads/Copy of TIBCO-NewPresentation-template-2026.pptx"
OUT      = r"c:/Work/github/flogo-enterprise-hub/demos/Agentic_AI/Power_Distribution_Use_Case/PowerDistribution_Architecture.pptx"

# ---- TIBCO 2026 brand palette (from template theme) ----
BLUE  = "3A8DDE"   # MCP
PURPLE= "6C40F3"   # Orchestrator
SLATE = "24323F"   # DB / dark text
GREEN = "1E9B6B"   # A2A
GRAY  = "B4B4B4"   # UI
MINT  = "B0DCCB"   # SMTP
LAV   = "CABDFB"   # LLM
WHITE = "FFFFFF"
TEXT  = "333333"
PANEL = "EEF2F6"   # light panel bg
CAP_ON_DARK = "DCE6F2"

def C(h): return RGBColor.from_string(h)

def layout_by_name(prs, name):
    for l in prs.slide_masters[0].slide_layouts:
        if l.name == name:
            return l
    raise KeyError(name)

def strip_slides(prs):
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rId = sid.get(qn('r:id'))
        if rId:
            prs.part.drop_rel(rId)   # unreference the part so it isn't serialized (no dup partnames)
        lst.remove(sid)

def set_title(slide, text, size=24, color=SLATE, bold=True):
    if slide.shapes.title is not None:
        tf = slide.shapes.title.text_frame
        tf.text = text
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = C(color)
        return slide.shapes.title

def _dim(v):  # positive integer EMU guard (Google Slides safe)
    e = int(Inches(v))
    return e if e > 0 else 1

def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_dim(x), _dim(y), _dim(w), _dim(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(2); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = C(col)
    return tb

def box(slide, x, y, w, h, lines, fill, title_col=WHITE, cap_col=CAP_ON_DARK,
        title_sz=11, cap_sz=7.5, border=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _dim(x), _dim(y), _dim(w), _dim(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    sh.line.color.rgb = C(border or fill); sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, (txt, is_title) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(1); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(title_sz if is_title else cap_sz)
        r.font.bold = bool(is_title)
        r.font.color.rgb = C(title_col if is_title else cap_col)
    return sh

def arrow(slide, x1, y1, x2, y2, color=SLATE, width=1.5):
    # nudge dead-straight lines so the connector bbox is never 0 wide/high (Google Slides safe)
    if abs(x2 - x1) < 0.01: x2 = x1 + 0.03
    if abs(y2 - y1) < 0.01: y2 = y1 + 0.03
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _dim(x1), _dim(y1), _dim(x2), _dim(y2))
    cn.line.color.rgb = C(color); cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    te = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(te)
    cn.shadow.inherit = False
    return cn

# =====================================================================
prs = Presentation(TEMPLATE)
strip_slides(prs)
L_TITLE = layout_by_name(prs, "TITLE")
L_CONTENT = layout_by_name(prs, "Title Only (Logo)")
L_SUMMARY = layout_by_name(prs, "Overview / Summary")

# ---------- S1: Title ----------
s = prs.slides.add_slide(L_TITLE)
if s.shapes.title is not None:
    s.shapes.title.text = "Electric Power Distribution — Residential Self-Service Assistant"
    for p in s.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = C(SLATE)
# subtitle placeholder (idx 1)
for ph in s.placeholders:
    if ph.placeholder_format.idx == 1:
        ph.text = ("A conversational AI that lets utility customers check bills, energy usage and "
                   "outages — and take action (report an outage, schedule a visit, request "
                   "reconnection) — over a real-time chat.")
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.color.rgb = C(TEXT)
textbox(s, 0.34, 0.42, 9.3, 0.35,
        [("AGENTIC AI ON TIBCO FLOGO ENTERPRISE", 12, True, PURPLE)], align=PP_ALIGN.CENTER)

# ---------- S2: The problem we automate ----------
s = prs.slides.add_slide(L_CONTENT)
set_title(s, "The problem we automate")
# TODAY panel
box(s, 0.45, 1.05, 4.35, 2.7, [("", False)], PANEL, border="D5DEE8")
textbox(s, 0.62, 1.18, 4.0, 2.5, [
    ("TODAY", 12, True, SLATE),
    ("• Call-center bound — routine billing, usage & outage questions tie up agents", 10.5, False, TEXT),
    ("• Low-risk actions (report outage, reconnect) handled manually", 10.5, False, TEXT),
    ("• After-hours gaps; storm-driven outage call spikes", 10.5, False, TEXT),
    ("• Inconsistent answers, little audit trail", 10.5, False, TEXT),
])
# WITH FLOGO panel
box(s, 5.05, 1.05, 4.5, 2.7, [("", False)], "EAF3FB", border="BFDCF3")
textbox(s, 5.22, 1.18, 4.16, 2.5, [
    ("WITH FLOGO AGENTIC AI", 12, True, BLUE),
    ("• 24/7 natural-language self-service", 10.5, False, TEXT),
    ("• Instant bill, usage & outage answers", 10.5, False, TEXT),
    ("• Guided actions with confirm-before-write guardrails", 10.5, False, TEXT),
    ("• Deflects call volume; consistent & auditable", 10.5, False, TEXT),
])
# metric strip
box(s, 0.45, 4.0, 9.1, 0.7,
    [("9 read tools   ·   4 action agents   ·   1 chat orchestrator   ·   all on TIBCO Flogo Enterprise", True)],
    PURPLE, title_sz=13)

# ---------- S3: Architecture ----------
s = prs.slides.add_slide(L_CONTENT)
set_title(s, "How it fits together — the 3-app pattern")
# boxes
box(s, 0.30, 2.35, 1.75, 0.95, [("Chatbot UI", True), ("(browser)", False)], GRAY, title_col=SLATE, cap_col=SLATE, title_sz=10.5, cap_sz=8)
box(s, 3.05, 1.02, 2.00, 0.52, [("OpenAI LLM — gpt-5.5", True)], LAV, title_col=SLATE, title_sz=9.5)
box(s, 2.55, 2.00, 2.55, 1.48, [
    ("AI Orchestrator", True),
    ("PowerDistributionAIOrchestrator", False),
    ("WebSocket :9680 · /grid", False),
    ("#wsserver → AI Agent → reply", False),
    ("intent routing + guardrails", False),
], PURPLE, title_sz=11, cap_sz=7.5)
box(s, 5.75, 1.30, 3.90, 0.95, [
    ("MCP Server — 9 read-only tools", True),
    ("PowerDistributionMCPServer · HTTP :9682 · /grid-bss", False),
], BLUE, title_sz=10.5, cap_sz=7.8)
box(s, 5.75, 2.55, 3.90, 0.95, [
    ("A2A Servers — 4 action agents", True),
    ("PowerDistributionA2AServers · :9683–9686", False),
], GREEN, title_sz=10.5, cap_sz=7.8)
box(s, 5.75, 4.10, 3.90, 0.72, [("PostgreSQL — power_distribution (10 tables)", True)], SLATE, title_sz=10)
box(s, 2.55, 4.10, 2.55, 0.72, [("SMTP / Email", True), ("confirmation e-mails", False)], MINT, title_col=SLATE, cap_col=SLATE, title_sz=10, cap_sz=7.5)
# arrows
arrow(s, 2.05, 2.80, 2.55, 2.74, SLATE)          # UI -> ORCH
arrow(s, 3.85, 2.00, 3.95, 1.55, PURPLE)          # ORCH -> LLM
arrow(s, 5.10, 2.45, 5.75, 1.90, BLUE)            # ORCH -> MCP
arrow(s, 5.10, 2.95, 5.75, 3.00, GREEN)           # ORCH -> A2A
arrow(s, 7.30, 2.25, 7.30, 4.10, BLUE)            # MCP -> DB
arrow(s, 7.90, 3.50, 7.90, 4.10, GREEN)           # A2A -> DB
arrow(s, 6.20, 3.50, 4.30, 4.10, GREEN)           # A2A -> SMTP
# WebSocket label + legend
textbox(s, 0.30, 3.45, 2.3, 0.3, [("WebSocket · :9680 · /grid", 8, True, PURPLE)], align=PP_ALIGN.CENTER)
textbox(s, 0.30, 5.02, 9.3, 0.32,
        [("MCP = reads   ·   A2A = writes   ·   Orchestrator = WebSocket brain", 9, True, SLATE)],
        align=PP_ALIGN.CENTER)

# ---------- S4: Capabilities ----------
s = prs.slides.add_slide(L_CONTENT)
set_title(s, "What the customer can do")
box(s, 0.45, 1.05, 4.55, 3.05, [("", False)], "EAF3FB", border="BFDCF3")
textbox(s, 0.62, 1.16, 4.2, 2.85, [
    ("ASK & LOOK UP   —   read-only · MCP", 11.5, True, BLUE),
    ("• Account & bill details", 10, False, TEXT),
    ("• Line-item bill explanation", 10, False, TEXT),
    ("• Energy usage & payment history", 10, False, TEXT),
    ("• Rate plans", 10, False, TEXT),
    ("• Current outages by area", 10, False, TEXT),
    ("• Ticket / appointment / request status", 10, False, TEXT),
])
box(s, 5.10, 1.05, 4.45, 3.05, [("", False)], "E9F6F0", border="BFE6D4")
textbox(s, 5.27, 1.16, 4.1, 2.85, [
    ("TAKE ACTION   —   guarded writes · A2A", 11.5, True, GREEN),
    ("• Report a power outage (opens a ticket + ETA)", 10, False, TEXT),
    ("• Schedule a field appointment", 10, False, TEXT),
    ("   (meter inspection/upgrade, tree trim, new service, survey)", 8.5, False, TEXT),
    ("• Submit reconnect / disconnect / transfer", 10, False, TEXT),
    ("• Email confirmation", 10, False, TEXT),
])
box(s, 0.45, 4.28, 9.1, 0.62,
    [("Every write is confirmed first — and reconnection is blocked while the account is past-due.", True)],
    SLATE, title_sz=11)

# ---------- S5: Walkthrough ----------
s = prs.slides.add_slide(L_CONTENT)
set_title(s, "Live walkthrough — two scenarios")
box(s, 0.45, 1.05, 4.55, 3.35, [("", False)], PANEL, border="D5DEE8")
textbox(s, 0.62, 1.16, 4.2, 3.15, [
    ("SCENARIO A · Report an outage", 11.5, True, PURPLE),
    ("1  “My power’s been out since 2 pm.”", 10, False, TEXT),
    ("2  Orchestrator confirms the service address", 10, False, TEXT),
    ("3  report_outage opens ticket OTKT-2026-#### (ETA ~4h)", 10, False, TEXT),
    ("4  “Email me the details” → send_confirmation_email", 10, False, TEXT),
])
box(s, 5.10, 1.05, 4.45, 3.35, [("", False)], PANEL, border="D5DEE8")
textbox(s, 5.27, 1.16, 4.1, 3.15, [
    ("SCENARIO B · Reconnect + guardrail", 11.5, True, GREEN),
    ("1  A disconnected customer asks to reconnect", 10, False, TEXT),
    ("2  Orchestrator checks the balance via MCP", 10, False, TEXT),
    ("3  If past-due → politely refuses, asks to pay first", 10, False, TEXT),
    ("4  Once cleared → submit_service_request (RECONNECT)", 10, False, TEXT),
    ("     → SRV-2026-#### submitted", 10, False, TEXT),
])
box(s, 0.45, 4.55, 9.1, 0.6,
    [("Reads → guarded write → confirmation, all in one natural conversation.", True)],
    PURPLE, title_sz=12)

# ---------- S6: Summary ----------
s = prs.slides.add_slide(L_SUMMARY)
set_title(s, "One pattern, any vertical")
textbox(s, 0.62, 1.0, 7.4, 3.0, [
    ("• Reads as MCP tools, writes as A2A agents, a WebSocket orchestrator as the brain — all on TIBCO Flogo Enterprise.", 12.5, False, TEXT),
    ("• Safe by design: reads vs. guarded writes · confirm-before-write · past-due reconnect guardrail.", 12.5, False, TEXT),
    ("• PostgreSQL-backed · port-isolated · each app independently deployable & scalable.", 12.5, False, TEXT),
    ("• Reusable pattern — the same wiring powers telecom, banking, hospital & airline demos; swap the domain, keep the architecture.", 12.5, False, TEXT),
    ("• For utilities: deflect call volume, 24/7 self-service, faster outage & reconnect handling.", 12.5, True, SLATE),
], anchor=MSO_ANCHOR.TOP)

prs.save(OUT)
print("SAVED:", OUT, "size(KB)=%.1f" % (os.path.getsize(OUT)/1024))
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
